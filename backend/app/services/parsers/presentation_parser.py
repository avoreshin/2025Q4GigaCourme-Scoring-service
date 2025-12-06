from app.services.parsers.base_parser import DocumentParser
from pptx import Presentation


class PresentationParser(DocumentParser):
    """Parser for PowerPoint presentations (.pptx)"""
    
    def parse(self, source: str) -> str:
        """Extract text from presentation slides"""
        text_parts = []
        
        try:
            prs = Presentation(source)
            for i, slide in enumerate(prs.slides, 1):
                slide_text = [f"Slide {i}:"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                if len(slide_text) > 1:  # More than just "Slide X:"
                    text_parts.append("\n".join(slide_text))
        except Exception as e:
            raise ValueError(f"Failed to parse presentation: {str(e)}")
        
        return "\n\n".join(text_parts)

